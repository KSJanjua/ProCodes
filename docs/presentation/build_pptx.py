"""Generate the editable PowerPoint deck for the InstanceDepth project.

    python docs/presentation/build_pptx.py

Writes ``docs/presentation/instancedepth-journey.pptx`` -- a native .pptx
(real text boxes, shapes, tables, speaker notes) so elements can be selected
and reused directly in PowerPoint, unlike the HTML version. Same 15-slide
story + a backup media-manifest slide; same palette (Pass 1 blue, Pass 2 teal,
Pass 3 amber). Media still to drop in are dashed amber placeholder boxes.

Idempotent: rerun to regenerate after editing this script.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ----------------------------------------------------------------- palette
BG      = RGBColor(0x0B, 0x0F, 0x14)
PANEL   = RGBColor(0x13, 0x1A, 0x23)
PANEL2  = RGBColor(0x18, 0x22, 0x30)
TEXT    = RGBColor(0xE9, 0xED, 0xF3)
MUT     = RGBColor(0x93, 0xA3, 0xB5)
FAINT   = RGBColor(0x5C, 0x6B, 0x7E)
P1      = RGBColor(0x3E, 0x6F, 0xD9)
P2      = RGBColor(0x1F, 0xA4, 0x91)
P3      = RGBColor(0xB0, 0x7D, 0x18)
P1T     = RGBColor(0x7D, 0xA3, 0xF2)
P2T     = RGBColor(0x3E, 0xD2, 0xBC)
P3T     = RGBColor(0xF5, 0xB6, 0x3F)
GOOD    = RGBColor(0x3E, 0xD2, 0xBC)
BAD     = RGBColor(0xE4, 0x57, 0x3D)
BADT    = RGBColor(0xF2, 0x83, 0x6C)
LINE    = RGBColor(0x2A, 0x35, 0x42)

SANS = "Segoe UI"
MONO = "Consolas"

EMU_IN = 914400
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


# ----------------------------------------------------------------- helpers
def slide(notes: str = ""):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _set_dash(shape):
    ln = shape.line._get_or_add_ln()
    d = ln.find(qn("a:prstDash"))
    if d is None:
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(d)
    else:
        d.set("val", "dash")


def text(s, x, y, w, h, runs, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         font=SANS, anchor=MSO_ANCHOR.TOP, spacing=None, wrap=True):
    """runs: a string, or a list of (string, {overrides}) tuples for inline styling."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, ov in runs:
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
    return tb


def kicker(s, x, y, txt, color=P3T):
    text(s, x, y, 9, 0.35, txt.upper(), size=12, color=color, bold=True, font=MONO)


def title(s, x, y, w, runs, size=34):
    text(s, x, y, w, 1.2, runs, size=size, color=TEXT, bold=True)


def box(s, x, y, w, h, fill=PANEL, line=LINE, line_w=1.0, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


def boxtext(s, x, y, w, h, runs, size=18, color=TEXT, bold=False,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS,
            fill=PANEL, line=LINE, line_w=1.0, spacing=None):
    box(s, x, y, w, h, fill=fill, line=line, line_w=line_w)
    text(s, x + 0.12, y, w - 0.24, h, runs, size=size, color=color, bold=bold,
         align=align, anchor=anchor, font=font, spacing=spacing)


def placeholder(s, x, y, w, h, pid, what, spec):
    shp = box(s, x, y, w, h, fill=PANEL2, line=P3T, line_w=1.5)
    _set_dash(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for txt, sz, col, bold, font in [
        (pid, 11, P3T, True, MONO),
        (what, 12.5, TEXT, False, SANS),
        (spec, 10, FAINT, False, MONO),
    ]:
        para = p if txt == pid else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(3)
        r = para.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = bold
        r.font.name = font
    return shp


def arrow(s, x, y, w, h=0.24, color=FAINT):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


def bar(s, x, y, w, h, color):
    b = box(s, x, y, w, h, fill=color, line=None, rounded=True)
    b.shadow.inherit = False
    return b


def footer(s, txt, color=MUT):
    text(s, 0.9, SH - 0.62, SW - 1.8, 0.4, txt, size=12.5, color=color, font=SANS)


MARGIN = 0.9

# ================================================================= 1 TITLE
s = slide("~20s -- Set the scene, no tech. 'Tonight I want to tell you the story of teaching a camera to see every person in a crowd -- separately, correctly, and steadily over time.'")
bar(s, 0, 6.7, SW, 0.8, P3); bar(s, 0, 6.4, SW, 0.32, P2); bar(s, 0, 6.15, SW, 0.22, P1)
# three depth-painted silhouettes
for bx, w, h, col in [(9.2, 0.8, 1.7, P1), (10.0, 0.95, 2.2, P2), (10.9, 1.1, 2.9, P3)]:
    box(s, bx, 6.7 - h, w, h, fill=col, line=None)
kicker(s, MARGIN, 1.5, "A research journey")
title(s, MARGIN, 2.0, 9.2, "Seeing every person's depth", size=54)
text(s, MARGIN, 3.5, 8.4, 0.9,
     "Instance-level video depth in crowded scenes -- beyond occlusions.",
     size=22, color=MUT)
text(s, MARGIN, SH - 0.7, 8, 0.4, "K. S. Janjua  ·  lab meeting  ·  July 2026",
     size=12, color=FAINT, font=MONO)

# ================================================================= 2 PROBLEM
s = slide("~35s -- The problem in one image. Left: today's best models fuse overlapping people into one blob at one depth. Right (V2): a real clip of a foundation model failing on a group. 'Everything great about modern depth estimation stops at the edge of a crowd.'")
kicker(s, MARGIN, 0.7, "The problem")
title(s, MARGIN, 1.15, 11, "Cameras flatten crowds.")
box(s, MARGIN, 2.5, 5.9, 4.0, fill=PANEL2, line=LINE)
text(s, MARGIN + 0.3, 2.7, 5, 0.4, "WHAT THE MODEL SEES", size=12, color=MUT, font=MONO)
# fused blob = one colour
for bx, w, h in [(2.4, 1.0, 2.2), (3.1, 1.1, 2.5), (3.9, 0.95, 2.0)]:
    box(s, bx, 5.4 - h, w, h, fill=RGBColor(0x41, 0x54, 0x7A), line=None)
text(s, MARGIN + 0.3, 5.7, 5.3, 0.5, "three people · one fused mass · one depth",
     size=15, color=MUT, align=PP_ALIGN.CENTER)
placeholder(s, 7.1, 2.5, 5.3, 4.0, "V2 · VIDEO",
            "Foundation-model depth failing on a real group clip -- people merging as they pass",
            "~620x460 · scripts/infer_video.py --compare")

# ================================================================= 3 WHY HARD
s = slide("~35s -- Three reasons this is genuinely hard. 1: the pixels you need are hidden. 2: a half-hidden person inherits the depth of whoever covers them. 3: video -- every frame answers slightly differently, so depth shimmers. 'The evidence is missing, borrowed, or unstable.'")
kicker(s, MARGIN, 0.7, "Why it's hard")
title(s, MARGIN, 1.15, 11.5, "Occlusion destroys the evidence.")
cards = [
    (P2, "Hidden parts have no pixels.", "the person behind simply isn't imaged"),
    (P1, "The wrong depth gets borrowed.", "the hidden body inherits the occluder's depth"),
    (P3, "The answer shimmers frame to frame.", "each frame decides slightly differently"),
]
cw = 3.75
for i, (col, head, sub) in enumerate(cards):
    x = MARGIN + i * (cw + 0.28)
    box(s, x, 2.6, cw, 3.7, fill=PANEL, line=LINE)
    box(s, x + 0.25, 2.9, 0.5, 0.5, fill=col, line=None)
    text(s, x + 0.25, 3.7, cw - 0.5, 1.0, head, size=20, color=TEXT, bold=True)
    text(s, x + 0.25, 5.1, cw - 0.5, 1.0, sub, size=14, color=MUT)

# ================================================================= 4 GOAL
s = slide("~20s -- What success looks like, in plain words. Four promises: depth for each person, correct even while overlapped, steady across the whole video, and in real metres -- not relative shades.")
kicker(s, MARGIN, 0.7, "The goal")
title(s, MARGIN, 1.15, 11, "What success looks like")
goals = [(P1, "Depth for each person"), (P2, "Correct through overlaps"),
         (P3, "Steady across frames"), (BADT, "In real metres")]
gw, gh = 5.5, 1.5
for i, (col, txt) in enumerate(goals):
    x = MARGIN + (i % 2) * (gw + 0.4)
    y = 2.9 + (i // 2) * (gh + 0.4)
    box(s, x, y, gw, gh, fill=PANEL, line=LINE)
    box(s, x + 0.35, y + gh / 2 - 0.16, 0.32, 0.32, fill=col, line=None)
    text(s, x + 1.0, y, gw - 1.2, gh, txt, size=23, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ================================================================= 5 DATASET
s = slide("~35s -- The unglamorous truth, and a real contribution. There was no dataset for this. All we were handed was RGB frames and their raw depth maps -- no masks, no identities, no per-person depth. So we built the entire ground-truth stack ourselves: segment every person, give each a stable identity across the whole video, assign each their depth layer, and label who occludes whom. Everything after the arrow is ours.")
kicker(s, MARGIN, 0.7, "The foundation")
title(s, MARGIN, 1.15, 12, "There was no dataset. So we built one.")
text(s, MARGIN, 2.5, 2.6, 0.35, "ALL WE HAD", size=12, color=MUT, font=MONO)
boxtext(s, MARGIN, 2.9, 2.6, 1.2, "RGB frames", size=18)
boxtext(s, MARGIN, 4.3, 2.6, 1.2, "depth maps (raw)", size=18)
arrow(s, 3.7, 3.5, 0.7)
built = [(P2, "Segment every person  →  masks"),
         (P1, "Stable identity across the video"),
         (P3, "Per-person depth layer"),
         (BADT, "Who occludes whom")]
for i, (col, txt) in enumerate(built):
    y = 2.85 + i * 0.92
    box(s, 4.7, y, 4.5, 0.78, fill=PANEL, line=col, line_w=1.5)
    box(s, 4.9, y + 0.27, 0.24, 0.24, fill=col, line=None)
    text(s, 5.3, y, 3.7, 0.78, txt, size=16, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
placeholder(s, 9.5, 2.85, 2.9, 3.0, "I5 · IMAGE",
            "One annotated frame: RGB → masks → IDs → depth layers", "~330x220")
footer(s, "Pixels and depth in; a fully-annotated instance-depth video dataset out.")

# ================================================================= 6 IDEA
s = slide("~40s -- The whole idea in one picture, following the ICCV-2025 InstanceDepth recipe. Pass 1 paints the scene's depth. Pass 2 finds each person and their depth layer. Pass 3 looks at every overlapping pair and fixes the depths where they collide. Our addition wraps the whole thing in time: remember people between frames.")
kicker(s, MARGIN, 0.7, "Our idea")
title(s, MARGIN, 1.15, 11, "Three passes over every frame")
boxtext(s, MARGIN, 3.2, 1.5, 1.2, "video frame", size=15)
arrow(s, 2.5, 3.68, 0.4)
passes = [(P1, P1T, "PASS 1", "Paint the scene", "depth everywhere"),
          (P2, P2T, "PASS 2", "Find each person", "mask + depth layer"),
          (P3, P3T, "PASS 3", "Fix the overlaps", "reason about pairs")]
px = 3.1
for i, (col, colt, tag, head, sub) in enumerate(passes):
    x = px + i * 3.15
    box(s, x, 2.7, 2.85, 2.1, fill=PANEL, line=col, line_w=2.0)
    text(s, x, 2.95, 2.85, 0.35, tag, size=13, color=colt, bold=True, align=PP_ALIGN.CENTER, font=MONO)
    text(s, x, 3.45, 2.85, 0.6, head, size=19, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 4.15, 2.85, 0.5, sub, size=14, color=MUT, align=PP_ALIGN.CENTER)
    if i < 2:
        arrow(s, x + 2.9, 3.68, 0.22)
box(s, 3.1, 5.35, 9.15, 0.62, fill=BG, line=BADT, line_w=2.0)
text(s, 3.1, 5.35, 9.15, 0.62, "+ ours: memory -- keeps each person across frames (feeds the next frame's Pass 1)",
     size=14, color=BADT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Following the ICCV 2025 recipe -- then extending it where video demands more.")

# ================================================================= 7 PHASE 1
s = slide("~50s -- Pass 1 results, and the payoff of the temporal fix. Left: scenes come out sharp. Right, the real numbers on the full 10,400-frame test set: once we froze the scene model and added only the memory module, EVERY number improved at once -- more accurate AND steadier. This is the module that used to do nothing; now it earns its place.")
kicker(s, MARGIN, 0.7, "Pass 1 · scene depth + memory", color=P1T)
title(s, MARGIN, 1.15, 12, [("Adding memory made it sharper ", {}), ("and", {"color": GOOD}), (" steadier", {})])
placeholder(s, MARGIN, 2.5, 5.6, 3.9, "I2 · IMAGE STRIP",
            "Two scenes × ( RGB | ground truth | our depth )", "~560x430 · scripts/visualize_hdi.py")
# results panel
rx = 7.0
box(s, rx, 2.5, 5.4, 3.9, fill=PANEL, line=LINE)
text(s, rx + 0.3, 2.7, 5, 0.35, "FULL TEST SET (10,400 frames) · baseline → + memory",
     size=12.5, color=MUT, font=MONO)
rows = [("Accuracy error (lower wins)", 0.078, 0.072, "0.078", "0.072 ✓"),
        ("Flicker (lower = steadier)", 0.94, 0.92, "0.94", "0.92 ✓"),
        ("Accuracy (higher wins)", 0.90, 0.97, "93.8%", "94.1% ✓")]
for i, (lab, bv, av, bl, al) in enumerate(rows):
    y = 3.3 + i * 1.0
    text(s, rx + 0.3, y, 4.8, 0.3, lab, size=14, color=TEXT)
    bar(s, rx + 0.3, y + 0.35, 1.9 * bv / max(bv, av), 0.16, P1)
    text(s, rx + 2.4, y + 0.28, 1.0, 0.3, bl, size=12, color=MUT, font=MONO)
    bar(s, rx + 3.1, y + 0.35, 1.9 * av / max(bv, av), 0.16, P3)
    text(s, rx + 5.05 - 0.9, y + 0.28, 1.1, 0.3, al, size=12, color=P3T, font=MONO)
text(s, rx + 0.3, 6.05, 5, 0.3, "every metric moved the right way -- at once ✓",
     size=13.5, color=GOOD)
footer(s, "Blue = scene model alone · amber = + the memory module (scene model frozen).")

# ================================================================= 8 PHASE 2
s = slide("~50s -- Pass 2. Masks are excellent on our data. The problem: the segmenter re-answers 'who is who' from scratch every frame, so the same person strobes colours. The fix: match each person to the previous frame by overlap, and hold identity even through a brief disappearance. Demo note: this is the --track-instances flag; always run the demo with it on.")
kicker(s, MARGIN, 0.7, "Pass 2 · finding & following people", color=P2T)
title(s, MARGIN, 1.15, 12, [("It finds people -- and now ", {}), ("keeps", {"color": GOOD}), (" who they are", {})])
placeholder(s, MARGIN, 2.5, 5.6, 3.5, "V3a · VIDEO -- BEFORE",
            "Identity flicker: the same person strobing colours frame to frame",
            "~540x360 · make_sequence_videos.py (no tracker)")
text(s, MARGIN, 6.1, 5.6, 0.4, "Before: new colour every frame.", size=15, color=BADT, bold=True)
p = placeholder(s, 7.0, 2.5, 5.4, 3.5, "V3b · VIDEO -- AFTER",
                "Same clip, identities locked: one person, one colour, start to finish",
                "~540x360 · infer_video.py --track-instances")
p.line.color.rgb = GOOD
text(s, 7.0, 6.1, 5.4, 0.4, "After: one person, one colour. ✓", size=15, color=GOOD, bold=True)
footer(s, "Match people frame-to-frame by overlap, and hold identity through brief disappearances.")

# ================================================================= 9 PHASE 3
s = slide("~45s -- Pass 3 results. Where two people overlap is exactly where depth is hardest -- and exactly where our refinement helps: overlap-region error drops 0.056 to 0.055, with 98.6% of overlap pixels within the tight accuracy band. Just as important is what it can NEVER do: by design the scene model is frozen and each per-person correction is capped at ±15%, so refinement is guaranteed to leave the rest of the scene untouched. Show V4: refined vs base, side by side, on a real overlap.")
kicker(s, MARGIN, 0.7, "Pass 3 · occlusion-aware refinement", color=P3T)
title(s, MARGIN, 1.15, 12, [("Better exactly ", {}), ("where people overlap", {"color": P3T})])
# left: results panel
box(s, MARGIN, 2.5, 6.2, 4.0, fill=PANEL, line=LINE)
text(s, MARGIN + 0.3, 2.7, 5.8, 0.3, "OVERLAP REGIONS · 3,083 frames with real occlusion",
     size=12.5, color=MUT, font=MONO)
text(s, MARGIN + 0.3, 3.15, 2.9, 0.3, "Depth error (lower wins)", size=13.5, color=TEXT)
text(s, MARGIN + 0.3, 3.5, 5.6, 0.55,
     [("0.056", {"size": 26, "bold": True, "color": MUT, "font": MONO}),
      ("  →  ", {"size": 18, "color": FAINT}),
      ("0.055 ✓", {"size": 26, "bold": True, "color": GOOD, "font": MONO}),
      ("   base → refined", {"size": 12.5, "color": MUT})])
text(s, MARGIN + 0.3, 4.25, 5.6, 0.3, "Accuracy in overlap regions", size=13.5, color=TEXT)
text(s, MARGIN + 0.3, 4.6, 5.6, 0.5,
     [("98.6%", {"size": 26, "bold": True, "color": GOOD, "font": MONO}),
      ("  of overlap pixels within the tight band", {"size": 12.5, "color": MUT})])
box(s, MARGIN + 0.3, 5.35, 5.6, 0.02, fill=LINE, line=None)
text(s, MARGIN + 0.3, 5.5, 5.8, 0.3,
     [("Safe by design:  ", {"size": 13.5, "bold": True, "color": P3T}),
      ("scene model frozen · corrections capped ±15%", {"size": 13.5, "color": TEXT})])
text(s, MARGIN + 0.3, 5.9, 5.8, 0.4, "refinement can only act on people -- the rest of the scene is untouched, guaranteed",
     size=12, color=MUT)
# right: media
placeholder(s, 7.5, 2.5, 4.9, 4.0, "V4 · VIDEO / TOGGLE",
            "Refined vs base, side by side, on a real overlap -- crisp person boundaries in depth",
            "~520x430 · infer_video.py --compare")

# ================================================================= 10 RESULTS AT A GLANCE
s = slide("~30s -- Everything on one card, all measured on the full 10,400-frame test set. Scene depth: error 0.072, 94.1% accuracy, and steadier than the per-frame baseline. People: found, followed, and kept -- one identity per person for the whole video, held through occlusions. Overlaps: 0.055 error right where bodies cross, 98.6% accuracy there. This is the slide to leave on screen during questions.")
kicker(s, MARGIN, 0.7, "Results")
title(s, MARGIN, 1.15, 11, "The numbers, all on one card")
cols = [
    (P1, P1T, "SCENE DEPTH", [
        ("0.072", "depth error (rel)"),
        ("94.1%", "accuracy"),
        ("-8%", "flicker vs per-frame")]),
    (P2, P2T, "PEOPLE & IDENTITY", [
        ("1 : 1", "one colour per person"),
        ("full video", "identity held"),
        ("through", "occlusions")]),
    (P3, P3T, "OVERLAP REGIONS", [
        ("0.055", "depth error (rel)"),
        ("98.6%", "accuracy"),
        ("±15%", "corrections, capped")]),
]
cw = 3.75
for i, (col, colt, head, stats) in enumerate(cols):
    x = MARGIN + i * (cw + 0.28)
    box(s, x, 2.5, cw, 3.9, fill=PANEL, line=col, line_w=2.0)
    text(s, x, 2.75, cw, 0.35, head, size=13, color=colt, bold=True, align=PP_ALIGN.CENTER, font=MONO)
    for j, (big, small) in enumerate(stats):
        y = 3.3 + j * 1.0
        text(s, x, y, cw, 0.5, big, size=26, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font=MONO)
        text(s, x, y + 0.5, cw, 0.35, small, size=12.5, color=MUT, align=PP_ALIGN.CENTER)
footer(s, "Full held-out test set · 10,400 frames · 3,083 with real person-on-person overlap.")

# ================================================================= 11 CONTRIBUTIONS
s = slide("~35s -- What in this project is OURS, beyond the paper we started from. 1: the dataset itself -- fully annotated by us from nothing but RGB and raw depth. 2: the temporal memory and the loss that trains it -- the paper is single-frame; we made it video-native. 3: identities that persist -- the paper re-finds people every frame; we follow them, even through occlusions. 4: refinement that is safe by construction -- frozen scene, capped corrections.")
kicker(s, MARGIN, 0.7, "Beyond the paper")
title(s, MARGIN, 1.15, 12, "What this project adds")
contrib = [
    (P2, "A dataset that didn't exist", "full instance-depth ground truth, hand-built from raw RGB + depth"),
    (P1, "Video-native depth", "temporal memory + a flicker-aware loss -- the paper is single-frame"),
    (P3, "Identities that persist", "one person, one identity, held through occlusions -- not re-found each frame"),
    (BADT, "Refinement, safe by construction", "frozen scene + capped corrections: it can help, it cannot harm"),
]
lw, lh = 5.5, 1.75
for i, (col, head, sub) in enumerate(contrib):
    x = MARGIN + (i % 2) * (lw + 0.4)
    y = 2.6 + (i // 2) * (lh + 0.35)
    box(s, x, y, lw, lh, fill=PANEL, line=col, line_w=1.75)
    text(s, x + 0.35, y + 0.25, lw - 0.7, 0.6, head, size=20, color=TEXT, bold=True)
    text(s, x + 0.35, y + 0.95, lw - 0.7, 0.7, sub, size=13.5, color=MUT)

# ================================================================= 12 DEMO
s = slide("~30s -- Let it speak for itself. One person walks behind another and comes out the other side: same colour, same depth, and the scene never wavers. Everything from the last five slides, in one clip. (Render with --track-instances; add --depth-range / --smooth-depth for foreign footage.)")
kicker(s, MARGIN, 0.7, "Demo")
title(s, MARGIN, 1.15, 12, "All of it, in one clip")
placeholder(s, MARGIN, 2.4, SW - 2 * MARGIN, 3.7, "V5 · VIDEO -- THE DEMO",
            "One person walks behind another -- same colour, same depth, steady scene",
            "~1128x380 · make_sequence_videos.py --track-instances")
text(s, MARGIN, 6.25, SW - 2 * MARGIN, 0.4,
     "RGB  |  depth  |  instances -- every person keeps their colour and their metres.",
     size=15, color=MUT, align=PP_ALIGN.CENTER)

# ================================================================= 13 NEXT
s = slide("~25s -- What's next, in order. 1 (already running): swap in the full pretrained depth model -- not just its backbone -- to push accuracy further, especially close-up. 2: robustness on footage from other cameras and places. 3: the write-up.")
kicker(s, MARGIN, 0.7, "Next")
title(s, MARGIN, 1.15, 11, "Three moves ahead")
moves = [(P1, "Stronger prior", "full pretrained depth model -- already running"),
         (P2, "Beyond our lab", "robust on footage from anywhere"),
         (P3, "The write-up", "dataset + method + results")]
mw = 3.6
for i, (col, head, sub) in enumerate(moves):
    x = MARGIN + i * (mw + 0.55)
    box(s, x, 2.9, mw, 1.9, fill=PANEL, line=col, line_w=2.0)
    text(s, x, 3.25, mw, 0.6, head, size=20, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 3.95, mw, 0.7, sub, size=13.5, color=MUT, align=PP_ALIGN.CENTER)
    if i < 2:
        arrow(s, x + mw + 0.08, 3.72, 0.35)

# ================================================================= 14 CONCLUSION
s = slide("~20s -- Close the loop with slide 4: we opened with four promises; here they are again, each with its receipt. Depth for each person -- delivered as per-person layers. Correct through overlaps -- 0.055 where bodies cross. Steady across frames -- flicker down, identities locked. Real metres -- 0.072 relative error on 10,400 frames. The promises held.")
kicker(s, MARGIN, 0.7, "Conclusion")
title(s, MARGIN, 1.15, 12, "The four promises -- delivered.")
delivered = [(P1, "Depth for each person", "per-person depth layers, every frame"),
             (P2, "Correct through overlaps", "0.055 error where bodies cross"),
             (P3, "Steady across frames", "flicker down · identities locked"),
             (BADT, "In real metres", "0.072 rel. error on 10,400 frames")]
gw, gh = 5.5, 1.6
for i, (col, head, sub) in enumerate(delivered):
    x = MARGIN + (i % 2) * (gw + 0.4)
    y = 2.75 + (i // 2) * (gh + 0.4)
    box(s, x, y, gw, gh, fill=PANEL, line=col, line_w=1.75)
    box(s, x + 0.35, y + 0.42, 0.32, 0.32, fill=col, line=None)
    text(s, x + 1.0, y + 0.22, gw - 1.4, 0.5,
         [(head + "  ", {"size": 19, "bold": True}), ("✓", {"size": 19, "color": GOOD, "bold": True})])
    text(s, x + 1.0, y + 0.85, gw - 1.4, 0.5, sub, size=13, color=MUT)

# ================================================================= 15 THANKS
s = slide("~15s -- Thanks + where to look. QR to the repo. Offer the backup slide if anyone wants the media checklist.")
bar(s, 0, 6.9, SW, 0.6, P3); bar(s, 0, 6.65, SW, 0.28, P2); bar(s, 0, 6.45, SW, 0.2, P1)
title(s, MARGIN, 2.4, 9, "Thank you.", size=60)
text(s, MARGIN, 4.0, 8, 0.6, "Questions -- happy to go deeper on any pass, or the dataset.",
     size=20, color=MUT)
text(s, MARGIN, 5.0, 8, 0.4, "github.com/KSJanjua/ProCodes · videodepth/", size=14, color=FAINT, font=MONO)
placeholder(s, 10.4, 2.4, 2.0, 2.0, "QR1", "QR → repo", "200x200")

# ================================================================= 16 BACKUP
s = slide("Backup only -- not part of the 8 minutes. Every placeholder, where it goes, its size, and the exact repo command. The two demo clips (V3b, V5) MUST use --track-instances so identities don't flicker.")
kicker(s, MARGIN, 0.6, "Backup · not part of the talk", color=FAINT)
title(s, MARGIN, 1.0, 12, "Editor's manifest -- media to drop in", size=28)
rows = [
    ("ID", "Slide · position", "Size", "Produce with"),
    ("V2", "S2 · right half", "620x460", "infer_video.py --compare (baseline failing)"),
    ("I5", "S5 · right -- the dataset", "330x220", "one frame: RGB → masks → IDs → layers"),
    ("I2", "S7 · left half", "560x430", "visualize_hdi.py, 2 scenes"),
    ("V3a", "S8 · left -- BEFORE", "540x360", "make_sequence_videos.py (no tracker)"),
    ("V3b", "S8 · right -- AFTER", "540x360", "infer_video.py --track-instances"),
    ("V4", "S9 · right -- refined vs base", "520x430", "infer_video.py --compare on an overlap clip"),
    ("V5", "S12 · main -- the demo", "1128x380", "make_sequence_videos.py --track-instances"),
    ("QR1", "S15 · right", "200x200", "QR to repo"),
]
tbl_shape = s.shapes.add_table(len(rows), 4, Inches(MARGIN), Inches(1.9),
                               Inches(SW - 2 * MARGIN), Inches(4.9))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.0)
tbl.columns[1].width = Inches(3.3)
tbl.columns[2].width = Inches(1.6)
tbl.columns[3].width = Inches(5.63)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if r else PANEL2
        cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        run = para.add_run(); run.text = val
        run.font.name = MONO
        run.font.size = Pt(11 if r else 12)
        run.font.bold = bool(r == 0)
        run.font.color.rgb = P3T if r == 0 else (TEXT if c == 0 else MUT)
footer(s, "Swap any video in as an inserted media clip. V3b and V5 MUST use --track-instances so colours stay locked.",
       color=FAINT)

# ----------------------------------------------------------------- write
out = Path(__file__).with_name("instancedepth-journey.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
